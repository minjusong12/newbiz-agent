from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt_from_outline(topic: str, slides: list, filename: str = None):
    # ✅ 회사 템플릿 적용 (같은 폴더에 company_template.pptx 가 있어야 함)
    prs = Presentation("company_template.pptx")

    # 템플릿에 맞는 레이아웃 인덱스 설정 (필요시 조정)
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    if not filename:
        filename = f"{topic}_슬라이드.pptx"

    # ✅ 1. 표지 슬라이드
    slide = prs.slides.add_slide(title_slide_layout)

    # 제목 처리
    if slide.shapes.title:
        slide.shapes.title.text = topic
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = txBox.text_frame
        tf.text = topic

    # 부제목/내용 자리 있으면 텍스트 입력
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "GPT 기반 자동 생성 자료"

    # ✅ 2. 본문 슬라이드 생성
    for s in slides:
        slide = prs.slides.add_slide(content_slide_layout)

        # 제목 처리
        if slide.shapes.title:
            slide.shapes.title.text = s["title"]
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = s["title"]

        # 본문 내용 작성
        content = s["summary"] + "\n\n📊 시각자료 제안: " + s["visual_suggestion"]
        if len(slide.placeholders) > 1:
            textbox = slide.placeholders[1]
            textbox.text = content

            for paragraph in textbox.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)

    prs.save(filename)
    print(f"✅ PPTX 저장 완료: {filename}")
